"""
Proposal export services for PDF, DOCX, and PPTX formats.
"""
from typing import List, Dict, Any, Optional
from pathlib import Path
from io import BytesIO
from datetime import datetime
import uuid

# PDF Export
try:
    from reportlab.lib.pagesizes import letter, A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.lib.colors import HexColor
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
    from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_JUSTIFY
    REPORTLAB_AVAILABLE = True
except ImportError:
    REPORTLAB_AVAILABLE = False

# DOCX Export
try:
    from docx import Document
    from docx.shared import Inches, Pt
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

# PPTX Export
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

class ProposalExporter:
    """Export proposals to various formats."""
    
    def __init__(self):
        self.export_dir = Path("./exports")
        self.export_dir.mkdir(exist_ok=True)
    
    def export_pdf(
        self,
        title: str,
        sections: List[Dict[str, Any]],
        project_name: str = None,
        client_name: str = None,
        company_name: str = None
    ) -> BytesIO:
        """
        Export proposal as modern, professionally styled PDF.
        
        Args:
            title: Proposal title
            sections: List of section dictionaries
            project_name: Optional project name
            client_name: Optional client name
        
        Returns:
            BytesIO buffer with PDF content
        """
        if not REPORTLAB_AVAILABLE:
            raise ValueError("ReportLab not available. Install with: pip install reportlab")
        
        buffer = BytesIO()
        doc = SimpleDocTemplate(
            buffer,
            pagesize=letter,
            rightMargin=50,
            leftMargin=50,
            topMargin=50,
            bottomMargin=50
        )
        
        # Modern color scheme
        primary_blue = HexColor('#1E40AF')  # Blue-800
        accent_blue = HexColor('#3B82F6')   # Blue-500
        dark_gray = HexColor('#1F2937')     # Gray-800
        medium_gray = HexColor('#6B7280')   # Gray-500
        light_gray = HexColor('#F3F4F6')    # Gray-100
        
        # Custom styles
        styles = getSampleStyleSheet()
        
        # Company/Header style
        company_style = ParagraphStyle(
            'CompanyName',
            parent=styles['Heading1'],
            fontSize=20,
            textColor=dark_gray,
            spaceAfter=6,
            fontName='Helvetica-Bold'
        )
        
        # Main title style
        title_style = ParagraphStyle(
            'MainTitle',
            parent=styles['Heading1'],
            fontSize=28,
            textColor=primary_blue,
            spaceAfter=10,
            alignment=TA_LEFT,
            fontName='Helvetica-Bold',
            leading=34
        )
        
        # Subtitle style
        subtitle_style = ParagraphStyle(
            'Subtitle',
            parent=styles['Normal'],
            fontSize=14,
            textColor=dark_gray,
            spaceAfter=20
        )
        
        # Section heading style
        section_heading_style = ParagraphStyle(
            'SectionHeading',
            parent=styles['Heading2'],
            fontSize=18,
            textColor=primary_blue,
            spaceAfter=12,
            spaceBefore=20,
            fontName='Helvetica-Bold'
        )
        
        # Body text style
        body_style = ParagraphStyle(
            'ModernBody',
            parent=styles['Normal'],
            fontSize=11,
            textColor=dark_gray,
            spaceAfter=12,
            leading=16,
            alignment=TA_JUSTIFY
        )
        
        # Info box style
        info_style = ParagraphStyle(
            'InfoBox',
            parent=styles['Normal'],
            fontSize=10,
            textColor=dark_gray,
            leftIndent=10,
            spaceAfter=6
        )
        
        # Build story
        story = []
        
        # Header section - use company_name if provided, otherwise fallback to "NovaIntel AI"
        company_display_name = company_name or "NovaIntel AI"
        story.append(Paragraph(company_display_name, company_style))
        story.append(Paragraph("<font color='#6B7280'>AI-Powered Proposal Platform</font>", styles['Normal']))
        story.append(Spacer(1, 0.3*inch))
        
        # Title
        story.append(Paragraph(f"<b>Request for Proposal</b>", title_style))
        story.append(Paragraph(f"Project: {title}", subtitle_style))
        story.append(Spacer(1, 0.2*inch))
        
        # Project details box
        details_html = f"""
        <para leftIndent='0' spaceAfter='8'>
        <b><font color='#1F2937'>RFP Details</font></b><br/>
        <font size='10' color='#6B7280'>────────────────────────────────</font>
        </para>
        """
        story.append(Paragraph(details_html, styles['Normal']))
        
        from utils.timezone import now_ist
        current_date = now_ist().strftime('%B %d, %Y')
        details_data = [
            ("Issue Date:", current_date),
            ("Client:", client_name or "[Client Name]"),
            ("Project:", project_name or title),
            ("Document ID:", f"RFP-{now_ist().strftime('%Y%m%d')}-{uuid.uuid4().hex[:6].upper()}"),
        ]
        
        for label, value in details_data:
            story.append(Paragraph(f"<b>{label}</b> {value}", info_style))
        
        story.append(Spacer(1, 0.3*inch))
        
        # Confidentiality notice
        confidentiality_text = f"""
        <para backColor='#FEF2F2' borderColor='#DC2626' borderWidth='1' borderPadding='10' borderRadius='5'>
        <b><font color='#991B1B'>Confidentiality Notice</font></b><br/>
        <font size='9' color='#7F1D1D'>
        This document contains confidential and proprietary information. 
        The contents may not be disclosed to any third party without express written consent. 
        All recipients are required to return or destroy this document upon request.
        </font>
        </para>
        """
        story.append(Paragraph(confidentiality_text, styles['Normal']))
        story.append(Spacer(1, 0.4*inch))
        
        # Sections
        section_counter = 1
        for section in sorted(sections, key=lambda x: x.get('order', 0)):
            section_title = section.get('title', '')
            section_content = section.get('content', '')
            
            if section_title:
                story.append(Paragraph(f"{section_counter}.0 {section_title}", section_heading_style))
                section_counter += 1
            
            if section_content:
                # Process content: handle bullet points, bold text, etc.
                content_lines = section_content.split('\n')
                for line in content_lines:
                    if line.strip():
                        # Handle bullet points
                        if line.strip().startswith('•') or line.strip().startswith('-'):
                            cleaned_line = line.strip().lstrip('•-').strip()
                            story.append(Paragraph(f"• {cleaned_line}", body_style))
                        elif line.strip().startswith(tuple(str(i) for i in range(10))):
                            # Numbered list
                            story.append(Paragraph(line.strip(), body_style))
                        else:
                            # Regular paragraph
                            story.append(Paragraph(line.strip(), body_style))
                    else:
                        story.append(Spacer(1, 0.1*inch))
                
                story.append(Spacer(1, 0.15*inch))
        
        # Footer
        story.append(Spacer(1, 0.3*inch))
        footer_text = f"""
        <para alignment='center' borderColor='#E5E7EB' borderWidth='1' borderPadding='10'>
        <font size='9' color='#6B7280'>
        Generated by {company_display_name} | {current_date}<br/>
        This is an AI-generated proposal based on your RFP analysis
        </font>
        </para>
        """
        story.append(Paragraph(footer_text, styles['Normal']))
        
        # Build PDF
        doc.build(story)
        buffer.seek(0)
        return buffer
    
    def export_docx(
        self,
        title: str,
        sections: List[Dict[str, Any]],
        project_name: str = None,
        client_name: str = None,
        company_name: str = None
    ) -> BytesIO:
        """
        Export proposal as modern, professionally styled DOCX.
        
        Args:
            title: Proposal title
            sections: List of section dictionaries
            project_name: Optional project name
            client_name: Optional client name
        
        Returns:
            BytesIO buffer with DOCX content
        """
        if not DOCX_AVAILABLE:
            raise ValueError("python-docx not available. Install with: pip install python-docx")
        
        from docx.shared import RGBColor
        from docx.oxml.ns import qn
        from docx.oxml import OxmlElement
        
        doc = Document()
        
        # Set document margins
        sections_obj = doc.sections
        for section_obj in sections_obj:
            section_obj.top_margin = Inches(0.7)
            section_obj.bottom_margin = Inches(0.7)
            section_obj.left_margin = Inches(0.7)
            section_obj.right_margin = Inches(0.7)
        
        # Helper function to add colored border to paragraph
        def add_bottom_line(paragraph, color='E5E7EB', width=4):
            """Add a subtle bottom line instead of full border"""
            p = paragraph._element
            pPr = p.get_or_add_pPr()
            pBdr = OxmlElement('w:pBdr')
            bottom = OxmlElement('w:bottom')
            bottom.set(qn('w:val'), 'single')
            bottom.set(qn('w:sz'), str(width))
            bottom.set(qn('w:space'), '1')
            bottom.set(qn('w:color'), color)
            pBdr.append(bottom)
            pPr.append(pBdr)
        
        # Helper function to add shading to paragraph
        def add_shading(paragraph, fill_color='FEF2F2'):
            shading_elm = OxmlElement('w:shd')
            shading_elm.set(qn('w:fill'), fill_color)
            paragraph._element.get_or_add_pPr().append(shading_elm)
        
        # Company header - use company_name if provided, otherwise fallback to "NovaIntel AI"
        company_display_name = company_name or "NovaIntel AI"
        company_para = doc.add_paragraph()
        run = company_para.add_run(company_display_name)
        run.bold = True
        run.font.size = Pt(20)
        run.font.color.rgb = RGBColor(31, 41, 55)  # Gray-800
        
        tagline = doc.add_paragraph('AI-Powered Proposal Platform')
        tagline.runs[0].font.size = Pt(10)
        tagline.runs[0].font.color.rgb = RGBColor(107, 114, 128)  # Gray-500
        
        doc.add_paragraph()  # Spacer
        
        # Main title with border
        title_para = doc.add_paragraph()
        title_run = title_para.add_run('Request for Proposal')
        title_run.bold = True
        title_run.font.size = Pt(28)
        title_run.font.color.rgb = RGBColor(30, 64, 175)  # Blue-800
        
        # Subtitle
        subtitle_para = doc.add_paragraph(f'Project: {title}')
        subtitle_para.runs[0].font.size = Pt(14)
        subtitle_para.runs[0].font.color.rgb = RGBColor(31, 41, 55)
        
        doc.add_paragraph()  # Spacer
        
        # RFP Details section
        details_heading = doc.add_paragraph('RFP Details')
        details_heading.runs[0].bold = True
        details_heading.runs[0].font.size = Pt(12)
        details_heading.runs[0].font.color.rgb = RGBColor(31, 41, 55)
        add_bottom_line(details_heading, 'E5E7EB', 4)
        
        from utils.timezone import now_ist
        current_date = now_ist().strftime('%B %d, %Y')
        from utils.timezone import now_ist
        doc_id = f"RFP-{now_ist().strftime('%Y%m%d')}-{uuid.uuid4().hex[:6].upper()}"
        
        details_data = [
            ('Issue Date:', current_date),
            ('Client:', client_name or '[Client Name]'),
            ('Project:', project_name or title),
            ('Document ID:', doc_id),
        ]
        
        for label, value in details_data:
            detail_para = doc.add_paragraph()
            label_run = detail_para.add_run(f'{label} ')
            label_run.bold = True
            label_run.font.size = Pt(10)
            value_run = detail_para.add_run(value)
            value_run.font.size = Pt(10)
            detail_para.paragraph_format.space_after = Pt(3)
        
        doc.add_paragraph()  # Spacer
        
        # Confidentiality notice with red shading
        conf_heading = doc.add_paragraph('Confidentiality Notice')
        conf_heading.runs[0].bold = True
        conf_heading.runs[0].font.color.rgb = RGBColor(153, 27, 27)  # Red-800
        conf_heading.runs[0].font.size = Pt(11)
        add_shading(conf_heading, 'FEF2F2')
        
        conf_text = doc.add_paragraph(
            'This document contains confidential and proprietary information. '
            'The contents may not be disclosed to any third party without express written consent. '
            'All recipients are required to return or destroy this document upon request.'
        )
        conf_text.runs[0].font.size = Pt(9)
        conf_text.runs[0].font.color.rgb = RGBColor(127, 29, 29)  # Red-900
        add_shading(conf_text, 'FEF2F2')
        conf_text.paragraph_format.space_after = Pt(6)
        
        doc.add_paragraph()  # Spacer
        doc.add_paragraph()  # Spacer
        
        # Sections
        section_counter = 1
        for section in sorted(sections, key=lambda x: x.get('order', 0)):
            section_title = section.get('title', '')
            section_content = section.get('content', '')
            
            if section_title:
                # Section heading without border
                heading_para = doc.add_heading(f"{section_counter}.0 {section_title}", level=1)
                heading_para.runs[0].font.color.rgb = RGBColor(30, 64, 175)  # Blue-800
                heading_para.runs[0].font.size = Pt(18)
                section_counter += 1
            
            if section_content:
                # Process content: handle bullet points, formatting
                content_lines = section_content.split('\n')
                for line in content_lines:
                    if line.strip():
                        para = doc.add_paragraph()
                        
                        # Handle bullet points
                        if line.strip().startswith('•') or line.strip().startswith('-'):
                            cleaned_line = line.strip().lstrip('•-').strip()
                            para.style = 'List Bullet'
                            para.add_run(cleaned_line)
                        elif line.strip().startswith(tuple(str(i) for i in range(10))):
                            # Numbered list
                            para.style = 'List Number'
                            para.add_run(line.strip())
                        else:
                            # Regular paragraph
                            para.add_run(line.strip())
                        
                        para.runs[0].font.size = Pt(11)
                        para.runs[0].font.color.rgb = RGBColor(31, 41, 55)
                        para.paragraph_format.space_after = Pt(8)
                        para.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
                    else:
                        doc.add_paragraph()  # Empty line for spacing
        
        # Footer
        doc.add_paragraph()  # Spacer
        footer_para = doc.add_paragraph()
        footer_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        footer_run = footer_para.add_run(
            f'Generated by {company_display_name} | {current_date}\n'
            'This is an AI-generated proposal based on your RFP analysis'
        )
        footer_run.font.size = Pt(9)
        footer_run.font.color.rgb = RGBColor(107, 114, 128)
        add_bottom_line(footer_para, 'E5E7EB', 4)
        
        # Save to buffer
        buffer = BytesIO()
        doc.save(buffer)
        buffer.seek(0)
        return buffer
    
    def export_pptx(
        self,
        title: str,
        sections: List[Dict[str, Any]],
        project_name: str = None,
        client_name: str = None,
        company_name: str = None
    ) -> BytesIO:
        """
        Export proposal as PowerPoint.
        
        Args:
            title: Proposal title
            sections: List of section dictionaries
            project_name: Optional project name
            client_name: Optional client name
        
        Returns:
            BytesIO buffer with PPTX content
        """
        if not PPTX_AVAILABLE:
            raise ValueError("python-pptx not available. Install with: pip install python-pptx")
        
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]
        
        title_shape.text = title
        subtitle_text = ""
        if client_name:
            subtitle_text += f"Client: {client_name}\n"
        if project_name:
            subtitle_text += f"Project: {project_name}\n"
        from utils.timezone import now_ist
        subtitle_text += now_ist().strftime('%B %d, %Y')
        subtitle_shape.text = subtitle_text
        
        # Section slides
        for section in sorted(sections, key=lambda x: x.get('order', 0)):
            section_title = section.get('title', '')
            section_content = section.get('content', '')
            
            if not section_title and not section_content:
                continue
            
            # Use title and content layout
            bullet_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide_layout)
            shapes = slide.shapes
            
            title_shape = shapes.title
            body_shape = shapes.placeholders[1]
            
            title_shape.text = section_title or "Section"
            
            # Add content as bullet points
            tf = body_shape.text_frame
            tf.text = section_content.split('\n')[0] if section_content else ""
            
            # Add remaining lines as bullet points
            for line in section_content.split('\n')[1:]:
                if line.strip():
                    p = tf.add_paragraph()
                    p.text = line.strip()
                    p.level = 0
        
        # Save to buffer
        buffer = BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        return buffer
    
    def save_export(
        self,
        buffer: BytesIO,
        format: str,
        proposal_id: int
    ) -> str:
        """
        Save export to file.
        
        Args:
            buffer: File buffer
            format: File format (pdf, docx, pptx)
            proposal_id: Proposal ID
        
        Returns:
            File path
        """
        filename = f"proposal_{proposal_id}_{uuid.uuid4().hex[:8]}.{format}"
        file_path = self.export_dir / filename
        
        with open(file_path, 'wb') as f:
            f.write(buffer.getvalue())
        
        return str(file_path)

# Global instance
proposal_exporter = ProposalExporter()

